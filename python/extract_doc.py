#!/usr/bin/env python3
"""
Extract plain text from a document file.

Usage:
    extract_doc.py <path> [--password PW]

Writes text to stdout. Status + warnings go to stderr. Exit codes
distinguish cases the caller needs to act on:

    0   ok (text written to stdout)
    2   usage / missing file
    3   unsupported extension
    4   generic extraction failure
    65  PDF is encrypted and no password was supplied
    66  PDF password is wrong

Supported formats: .docx, .xlsx, .pptx, .pdf, and the text-ish
formats (.md .txt .csv) — though those are better handled in Node.
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path


def _docx(path: Path) -> str:
    from docx import Document  # python-docx

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"# {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells))
        parts.append("")
    return "\n".join(parts).rstrip()


def _pptx(path: Path) -> str:
    from pptx import Presentation  # python-pptx

    pres = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(pres.slides, start=1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            # Tables inside slides.
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        parts.append("")
    return "\n".join(parts).rstrip()


def _pdf(path: Path, password: str | None) -> str:
    """Extract PDF text with pypdf. Handles encryption: if the PDF is
    encrypted and no password was supplied, exit 65 so the Node side can
    prompt; if the supplied password is wrong, exit 66."""
    from pypdf import PdfReader
    from pypdf.errors import DependencyError, FileNotDecryptedError

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        if not password:
            print("PDF_ENCRYPTED", file=sys.stderr)
            sys.exit(65)
        try:
            result = reader.decrypt(password)
        except DependencyError as e:
            print(f"PDF decryption needs an extra dep: {e}", file=sys.stderr)
            sys.exit(4)
        # pypdf.decrypt returns 0 (failure), 1 (user pw), 2 (owner pw)
        if result == 0:
            print("PDF_BAD_PASSWORD", file=sys.stderr)
            sys.exit(66)

    parts: list[str] = []
    try:
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
    except FileNotDecryptedError:
        print("PDF_BAD_PASSWORD", file=sys.stderr)
        sys.exit(66)
    return "\n\n".join(parts)


def _textish(path: Path) -> str:
    """Treat .md / .txt / .csv as UTF-8 with a permissive fallback."""
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="utf-8", errors="replace")


_HANDLERS = {
    ".docx": _docx,
    ".xlsx": _xlsx,
    ".pptx": _pptx,
    ".md": _textish,
    ".txt": _textish,
    ".csv": _textish,
}


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("path", type=Path)
    ap.add_argument("--password", default=None)
    args = ap.parse_args()

    if not args.path.exists():
        print(f"not found: {args.path}", file=sys.stderr)
        return 2
    ext = args.path.suffix.lower()

    try:
        if ext == ".pdf":
            text = _pdf(args.path, args.password)
        elif ext in _HANDLERS:
            text = _HANDLERS[ext](args.path)
        else:
            print(f"unsupported extension: {ext}", file=sys.stderr)
            return 3
    except SystemExit:
        raise
    except Exception as e:
        print(f"extract failed: {e}", file=sys.stderr)
        return 4

    sys.stdout.write(text)
    return 0


if __name__ == "__main__":
    sys.exit(main())
